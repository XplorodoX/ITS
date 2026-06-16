# /// script
# dependencies = [
#     "python-pptx",
# ]
# ///

"""
Generates a professional, factual widescreen PowerPoint presentation for the AALeC Quiz system.
Theme: Clean Dark (Deep Slate, Dark Slate Gray cards, Clean White & Slate text, Standard Blue/Green accents)
Usage:
    uv run generate_pptx.py
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# ── COLOR PALETTE (Clean, high contrast, professional) ───────────────────────
BG_COLOR = RGBColor(15, 18, 25)       # #0f1219 - Deep dark slate background
CARD_COLOR = RGBColor(24, 28, 37)     # #181c25 - Slate gray cards
BORDER_COLOR = RGBColor(38, 45, 58)   # #262d3a - Card borders
TEXT_TITLE = RGBColor(56, 189, 248)   # #38bdf8 - Slate Cyan
TEXT_WHITE = RGBColor(241, 245, 249)  # #f1f5f9 - Off-white main text
TEXT_MUTED = RGBColor(148, 163, 184)  # #94a3b8 - Muted slate text
ACCENT_GREEN = RGBColor(34, 197, 94)  # #22c55e - Success Green
ACCENT_BLUE = RGBColor(37, 99, 235)   # #2563eb - Primary Blue
ACCENT_RED = RGBColor(239, 68, 68)    # #ef4444 - Alert Red
ACCENT_YELLOW = RGBColor(234, 179, 8)  # #eab308 - Warning Yellow

FONT_TITLE = "Arial"
FONT_BODY = "Calibri"

def set_slide_bg_and_border(prs, slide):
    """Sets a solid dark background on a slide by adding a background rectangle."""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_COLOR
    bg.line.color.rgb = BG_COLOR # Hide border
    
    # Send background to back
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return bg

def add_header(slide, title_text):
    """Adds a standard header to the slide."""
    header_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(11.83), Inches(0.8))
    tf = header_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
    
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = FONT_TITLE
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE

def create_card(slide, left, top, width, height, bg_color=CARD_COLOR, border_color=BORDER_COLOR):
    """Adds a card shape to organize content visually."""
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    card.line.color.rgb = border_color
    card.line.width = Pt(1.0)
    return card

def add_bullet_points(tf, points, text_color=TEXT_MUTED, font_size=16, spacing=8):
    """Adds formatted bullet points to a text frame."""
    for i, pt in enumerate(points):
        if i == 0 and not tf.paragraphs[0].text:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = pt
        p.font.name = FONT_BODY
        p.font.size = Pt(font_size)
        p.font.color.rgb = text_color
        p.space_after = Pt(spacing)
        p.level = 0

def add_footer_page_number(slide, num, total):
    """Adds slide page numbers at the bottom right."""
    footer_box = slide.shapes.add_textbox(Inches(11.0), Inches(6.9), Inches(1.58), Inches(0.3))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num} / {total}"
    p.alignment = PP_ALIGN.RIGHT
    p.font.name = FONT_BODY
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_MUTED

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_slide_layout = prs.slide_layouts[6]
    total_slides = 9
    
    # ── SLIDE 1: TITLE SLIDE ──────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg_and_border(prs, slide)
    
    # Main Title
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(6.5), Inches(3.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    p = tf.paragraphs[0]
    p.text = "AALeC Quiz\nSystem"
    p.font.name = FONT_TITLE
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(12)
    
    p2 = tf.add_paragraph()
    p2.text = "Technische Dokumentation des IoT-Quiz-Systems mit\nESP8266-Hardware, Python-Backend und Next.js-Beamer-Schnittstelle."
    p2.font.name = FONT_BODY
    p2.font.size = Pt(18)
    p2.font.color.rgb = TEXT_MUTED
    
    # Tech tags (small container box)
    tags_box = slide.shapes.add_textbox(Inches(0.75), Inches(5.0), Inches(6.5), Inches(1.0))
    tf_tags = tags_box.text_frame
    tf_tags.word_wrap = True
    tf_tags.margin_left = tf_tags.margin_right = tf_tags.margin_top = tf_tags.margin_bottom = 0
    p_tags = tf_tags.paragraphs[0]
    p_tags.text = "Hardware: ESP8266 (C++)  |  Backend: Python (paho-mqtt)\nBroker: Mosquitto  |  Frontend-Anzeige: Next.js (TypeScript)"
    p_tags.font.name = FONT_BODY
    p_tags.font.size = Pt(13)
    p_tags.font.bold = True
    p_tags.font.color.rgb = TEXT_TITLE
    
    # Embed Mockup Image
    img_path = os.path.join(os.path.dirname(__file__), "assets", "buzzer_mockup.png")
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(7.5), Inches(1.2), width=Inches(5.0), height=Inches(5.0))
        
    add_footer_page_number(slide, 1, total_slides)
    
    # ── SLIDE 2: SYSTEMDESCRIPTION ─────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg_and_border(prs, slide)
    add_header(slide, "Systembeschreibung")
    
    # Left Column: Bullet points
    text_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(5.8), Inches(4.8))
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    add_bullet_points(tf, [
        "Das AALeC-Quiz ist ein hardwarebasiertes Multiplayer-Quizsystem für Gruppenaktivitäten.",
        "Spieler nutzen physische Client-Geräte (AALeC V3) zur Eingabe von Antworten und Messwerten anstelle von Smartphones.",
        "Die Datenübertragung erfolgt ereignisgesteuert über ein lokales WLAN-Netzwerk mittels des MQTT-Protokolls.",
        "Die Beamer UI (Next.js-Präsentationsschicht) visualisiert Fragen, Timer und Ranglisten in Echtzeit."
    ], spacing=15)
    
    # Right Column: Visual highlight card
    create_card(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(4.8))
    info_box = slide.shapes.add_textbox(Inches(7.4), Inches(1.9), Inches(4.7), Inches(4.0))
    tf_info = info_box.text_frame
    tf_info.word_wrap = True
    tf_info.margin_left = tf_info.margin_right = tf_info.margin_top = tf_info.margin_bottom = 0
    
    p = tf_info.paragraphs[0]
    p.text = "Spielablauf"
    p.font.name = FONT_TITLE
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = TEXT_TITLE
    p.space_after = Pt(12)
    
    p2 = tf_info.add_paragraph()
    p2.text = "Die physischen Controller verbinden sich nach dem Start automatisch mit dem lokalen Access Point, registrieren sich am Game Master und warten auf den Spielstart.\n\n" \
              "Die Fragen-Metadaten werden parallel verteilt. Die Antwortdaten (Werte und Latenz in Millisekunden) werden über den Broker zurück an das Python-Backend übermittelt."
    p2.font.name = FONT_BODY
    p2.font.size = Pt(15)
    p2.font.color.rgb = TEXT_WHITE
    p2.space_after = Pt(10)
    
    add_footer_page_number(slide, 2, total_slides)
    
    # ── SLIDE 3: NETWORK & SYSTEM ARCHITECTURE ─────────────────────────────────
    slide = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg_and_border(prs, slide)
    add_header(slide, "Netzwerk- und Systemarchitektur")
    
    # Left Column: Explanation
    text_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(5.8), Inches(4.8))
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    add_bullet_points(tf, [
        "WLAN-Access-Point: Stellt das lokale Netzwerk 'AALeC-Quiz' (192.168.4.0/24) bereit.",
        "MQTT Broker (Mosquitto): Verwaltet Verbindungen für Clients (TCP 1883) und Beamer-Anzeige (WebSockets 9001).",
        "Game Master (Python): Führt das Regelsystem aus, verarbeitet Antworten und steuert Statusänderungen.",
        "Beamer UI (Next.js): Dashboard im Browser des Projektions-Rechners, holt Spieldaten live über MQTT-WebSockets."
    ], spacing=12)
    
    # Right Column: Visual Diagram Blocks using shapes
    diag_left = Inches(7.2)
    diag_top = Inches(1.5)
    
    # 1. Hotspot
    create_card(slide, diag_left + Inches(1.2), diag_top, Inches(2.5), Inches(1.0), bg_color=CARD_COLOR, border_color=BORDER_COLOR)
    box1 = slide.shapes.add_textbox(diag_left + Inches(1.3), diag_top + Inches(0.15), Inches(2.3), Inches(0.7))
    box1.text_frame.word_wrap = True
    p1 = box1.text_frame.paragraphs[0]
    p1.text = "WLAN AP\nSSID: AALeC-Quiz"
    p1.font.size = Pt(13)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_WHITE
    p1.alignment = PP_ALIGN.CENTER
    
    # Arrow Down 1
    arrow1 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, diag_left + Inches(2.2), diag_top + Inches(1.05), Inches(0.5), Inches(0.4))
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = BORDER_COLOR
    arrow1.line.color.rgb = BORDER_COLOR
    
    # 2. MQTT Broker
    create_card(slide, diag_left + Inches(1.2), diag_top + Inches(1.5), Inches(2.5), Inches(1.0), bg_color=CARD_COLOR, border_color=BORDER_COLOR)
    box2 = slide.shapes.add_textbox(diag_left + Inches(1.3), diag_top + Inches(1.65), Inches(2.3), Inches(0.7))
    box2.text_frame.word_wrap = True
    p2 = box2.text_frame.paragraphs[0]
    p2.text = "Mosquitto Broker\nPort 1883 & 9001 (WS)"
    p2.font.size = Pt(13)
    p2.font.bold = True
    p2.font.color.rgb = TEXT_WHITE
    p2.alignment = PP_ALIGN.CENTER
    
    # Left: AALeC Controllers
    create_card(slide, diag_left, diag_top + Inches(3.2), Inches(2.1), Inches(1.2), bg_color=CARD_COLOR, border_color=BORDER_COLOR)
    box3 = slide.shapes.add_textbox(diag_left + Inches(0.1), diag_top + Inches(3.3), Inches(1.9), Inches(1.0))
    box3.text_frame.word_wrap = True
    p3 = box3.text_frame.paragraphs[0]
    p3.text = "AALeC Controller\nESP8266 Client\n(TCP 1883)"
    p3.font.size = Pt(11)
    p3.font.bold = True
    p3.font.color.rgb = TEXT_WHITE
    p3.alignment = PP_ALIGN.CENTER
    
    # Right: Python GM & Next.js UI
    create_card(slide, diag_left + Inches(2.8), diag_top + Inches(3.2), Inches(2.3), Inches(1.2), bg_color=CARD_COLOR, border_color=BORDER_COLOR)
    box4 = slide.shapes.add_textbox(diag_left + Inches(2.9), diag_top + Inches(3.3), Inches(2.1), Inches(1.0))
    box4.text_frame.word_wrap = True
    p4 = box4.text_frame.paragraphs[0]
    p4.text = "Mac OS Host (Docker)\n- Python GM (1883)\n- Next.js UI (WS 9001)"
    p4.font.size = Pt(11)
    p4.font.bold = True
    p4.font.color.rgb = TEXT_WHITE
    p4.alignment = PP_ALIGN.CENTER
    
    # Connectors
    c1 = slide.shapes.add_shape(MSO_SHAPE.LEFT_RIGHT_ARROW, diag_left + Inches(2.15), diag_top + Inches(2.6), Inches(0.6), Inches(0.5))
    c1.fill.solid()
    c1.fill.fore_color.rgb = BORDER_COLOR
    c1.line.color.rgb = BORDER_COLOR
    
    c2 = slide.shapes.add_shape(MSO_SHAPE.LEFT_RIGHT_ARROW, diag_left + Inches(2.15), diag_top + Inches(3.55), Inches(0.6), Inches(0.5))
    c2.fill.solid()
    c2.fill.fore_color.rgb = BORDER_COLOR
    c2.line.color.rgb = BORDER_COLOR
    
    add_footer_page_number(slide, 3, total_slides)
    
    # ── SLIDE 4: CONTROLLER HARDWARE (AALeC V3) ────────────────────────────────
    slide = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg_and_border(prs, slide)
    add_header(slide, "Controller-Hardware (AALeC V3)")
    
    # Left Column: Specs
    text_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(5.8), Inches(4.8))
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    add_bullet_points(tf, [
        "Prozessor: ESP8266 NodeMCU mit integriertem WLAN-Modul.",
        "Display: Monochrome OLED-Anzeige (128x64 Pixel) zur Darstellung von Verbindungsinformationen und Antwortoptionen.",
        "LED-Statusleiste: 5x WS2812B adressierbare LEDs visualisieren den Zeitverlauf (Countdown) oder Systemzustände.",
        "Bedienelemente: Drehgeber (Rotary Encoder) mit Taster dient zur Menüeingabe und Antwortbestätigung.",
        "Sensorik: Analoges Potentiometer (0-100%) und DS18B20-Temperatursensor für physikalische Messaufgaben."
    ], spacing=10)
    
    # Right Column: Image
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(7.5), Inches(1.5), width=Inches(4.8), height=Inches(4.8))
        
    add_footer_page_number(slide, 4, total_slides)
    
    # ── SLIDE 5: BACKEND GAME MASTER ──────────────────────────────────────────
    slide = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg_and_border(prs, slide)
    add_header(slide, "Backend - Game Master")
    
    # Left Column: State machine & general tasks
    text_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(5.8), Inches(4.8))
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    add_bullet_points(tf, [
        "Python-Dienst zur logischen Steuerung des Quizablaufs.",
        "Kern-Zustandsmaschine: WAITING ➔ QUESTION ➔ VOTING ➔ REVEAL ➔ SCORES ➔ (NEXT/ENDED).",
        "Antwortüberwachung: Schließt Abstimmungsrunden vorzeitig, sobald alle registrierten Online-Spieler geantwortet haben.",
        "Verbindungsstatus: Erkennt Disconnects über LWT-Nachrichten (Last Will) und passt die aktive Spieleranzahl dynamisch an."
    ], spacing=12)
    
    # Right Column: Scoring card
    create_card(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(4.8))
    score_box = slide.shapes.add_textbox(Inches(7.3), Inches(1.7), Inches(4.9), Inches(4.4))
    tf_score = score_box.text_frame
    tf_score.word_wrap = True
    tf_score.margin_left = tf_score.margin_right = tf_score.margin_top = tf_score.margin_bottom = 0
    
    p = tf_score.paragraphs[0]
    p.text = "Punktemodelle"
    p.font.name = FONT_TITLE
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = TEXT_TITLE
    p.space_after = Pt(12)
    
    p2 = tf_score.add_paragraph()
    p2.text = "- MCQ & Höher/Niedriger:\n" \
              "  1000 Basispunkte bei korrekter Antwort + maximal 500 Zeitbonus (linear sinkend über die Antwortdauer) + Serie-Bonus (+200 pro Runde in Serie, maximal +600).\n\n" \
              "- Schätzfragen:\n" \
              "  Gestufte Punkteverteilung nach relativer Abweichung zum echten Wert (100% Punkte bei exakt, 80% innerhalb 5%, abnehmend bis 20% innerhalb 30% Abweichung).\n\n" \
              "- Poti- & Temperatur-Challenges:\n" \
              "  Lineare Punkteermittlung innerhalb des Toleranzfensters (z.B. ±5% oder ±1.5°C). Volle Punkte bei exakter Übereinstimmung."
    p2.font.name = FONT_BODY
    p2.font.size = Pt(12.5)
    p2.font.color.rgb = TEXT_WHITE
    
    add_footer_page_number(slide, 5, total_slides)
    
    # ── SLIDE 6: BEAMER UI ────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg_and_border(prs, slide)
    add_header(slide, "Beamer-Präsentationsschicht")
    
    # Left Column: Beamer UI details
    text_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(5.8), Inches(4.8))
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    add_bullet_points(tf, [
        "Web-Applikation auf Basis von Next.js, React und TypeScript.",
        "Kommunikation: Nutzt MQTT.js über WebSockets zur Echtzeit-Synchronisation mit dem Broker.",
        "Unterstützte Ansichten:\n" \
        "  - Lobby: Spielerauflistung und Bereitschaftsstatus. Ermöglicht Spielstart.\n" \
        "  - Frage: Darstellung der aktiven Aufgabe, Antwortoptionen und verbleibenden Ablaufzeit.\n" \
        "  - Auflösung (Reveal): Verteilung der abgegebenen Antworten.\n" \
        "  - Rangliste (Scoreboard): Platzierungen sortiert nach Gesamtpunkten."
    ], spacing=12)
    
    # Right Column: UI visualization card
    create_card(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(4.8))
    ui_box = slide.shapes.add_textbox(Inches(7.3), Inches(1.8), Inches(4.9), Inches(4.2))
    tf_ui = ui_box.text_frame
    tf_ui.word_wrap = True
    tf_ui.margin_left = tf_ui.margin_right = tf_ui.margin_top = tf_ui.margin_bottom = 0
    
    p = tf_ui.paragraphs[0]
    p.text = "Gestaltung und Anzeige"
    p.font.name = FONT_TITLE
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = TEXT_TITLE
    p.space_after = Pt(12)
    
    p2 = tf_ui.add_paragraph()
    p2.text = "Das User Interface ist speziell für die Projektion auf Leinwände optimiert. Kontrastreiche, einheitliche Zustandselemente sorgen für Übersichtlichkeit.\n\n" \
              "Die Rangliste verwendet automatische Sortierungen und Platzierungscharts, um Verschiebungen in der Punkteverteilung nach jeder Spielrunde abzubilden."
    p2.font.name = FONT_BODY
    p2.font.size = Pt(14)
    p2.font.color.rgb = TEXT_WHITE
    p2.space_after = Pt(10)
    
    add_footer_page_number(slide, 6, total_slides)
    
    # ── SLIDE 7: MQTT PROTOCOL ────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg_and_border(prs, slide)
    add_header(slide, "MQTT-Kommunikationsprotokoll")
    
    # Table like layout using shapes
    table_left = Inches(0.75)
    table_top = Inches(1.5)
    
    # Header row
    r_height = Inches(0.5)
    create_card(slide, table_left, table_top, Inches(11.83), r_height, bg_color=RGBColor(30, 35, 48))
    hdr_box = slide.shapes.add_textbox(table_left, table_top + Inches(0.05), Inches(11.83), r_height)
    tf_hdr = hdr_box.text_frame
    p = tf_hdr.paragraphs[0]
    p.text = "   Topic                                       Richtung              Bedeutung / Payload"
    p.font.name = FONT_TITLE
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = TEXT_TITLE
    
    topics = [
        ("quiz/state", "GM -> Alle", "Überträgt den aktuellen Systemzustand (WAITING, VOTING, etc.)"),
        ("quiz/question", "GM -> Alle", "Überträgt Fragendetails (Typ, Limits, Antwortoptionen)"),
        ("quiz/connect/+", "AALeC -> GM", "Anmeldung eines Controllers am Game Master"),
        ("quiz/disconnect/+", "AALeC -> GM", "Last-Will-Topic: Meldet Verbindungsabbruch eines Clients"),
        ("quiz/answer/+", "AALeC -> GM", "Absenden der Antwort inklusive gemessener Latenzzeit"),
        ("quiz/control", "Web -> GM", "Schnittstelle für Ablaufbefehle (starten, zurücksetzen)")
    ]
    
    for i, (topic, direct, desc) in enumerate(topics):
        curr_top = table_top + Inches(0.6) + Inches(i * 0.7)
        create_card(slide, table_left, curr_top, Inches(11.83), Inches(0.6))
        
        row_box = slide.shapes.add_textbox(table_left + Inches(0.2), curr_top + Inches(0.08), Inches(11.4), Inches(0.5))
        tf_row = row_box.text_frame
        p_row = tf_row.paragraphs[0]
        
        # Color coding direction
        dir_color = ACCENT_GREEN if "GM" in direct.split("->")[0] else ACCENT_BLUE
        
        p_row.text = f"{topic:<40} {direct:<20} {desc}"
        p_row.font.name = FONT_BODY
        p_row.font.size = Pt(13)
        p_row.font.color.rgb = TEXT_WHITE
        
    add_footer_page_number(slide, 7, total_slides)
    
    # ── SLIDE 8: DEVOPS & OPERATION ───────────────────────────────────────────
    slide = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg_and_border(prs, slide)
    add_header(slide, "Inbetriebnahme & Bereitstellung")
    
    # 4 Steps Horizontal Layout
    step_width = Inches(2.7)
    step_height = Inches(4.5)
    step_spacing = Inches(0.3)
    start_left = Inches(0.75)
    top_pos = Inches(1.6)
    
    steps = [
        ("SCHRITT 1", "WLAN-Setup", "Access Point bereitstellen:\n\nSSID: 'AALeC-Quiz'\nPasswort: '12345678'\nIP-Bereich: 192.168.4.X\n\nBroker-IP ermitteln.", ACCENT_BLUE),
        ("SCHRITT 2", "Firmware Flash", "Einrichten der Broker-IP in 'config.h'.\n\nKompilieren und flashen des Controllers über PlatformIO:\n'pio run -t upload'", ACCENT_YELLOW),
        ("SCHRITT 3", "Docker Compose", "Starten der Server-Infrastruktur auf dem Host-System:\n\n'docker compose up -d'\n\n(Mosquitto, Game Master, Web)", ACCENT_GREEN),
        ("SCHRITT 4", "Spielsteuerung", "Beamer-Dashboard im Webbrowser aufrufen.\n\nPhysische Controller verbinden sich automatisch nach dem Start.", ACCENT_RED)
    ]
    
    for i, (num, title, desc, color) in enumerate(steps):
        left_pos = start_left + Inches(i * 3.0)
        create_card(slide, left_pos, top_pos, step_width, step_height, border_color=BORDER_COLOR)
        
        t_box = slide.shapes.add_textbox(left_pos + Inches(0.15), top_pos + Inches(0.2), step_width - Inches(0.3), step_height - Inches(0.4))
        tf_step = t_box.text_frame
        tf_step.word_wrap = True
        
        p = tf_step.paragraphs[0]
        p.text = num
        p.font.name = FONT_TITLE
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = color
        p.space_after = Pt(8)
        
        p2 = tf_step.add_paragraph()
        p2.text = title
        p2.font.name = FONT_TITLE
        p2.font.size = Pt(16)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_WHITE
        p2.space_after = Pt(12)
        
        p3 = tf_step.add_paragraph()
        p3.text = desc
        p3.font.name = FONT_BODY
        p3.font.size = Pt(12)
        p3.font.color.rgb = TEXT_MUTED
        
    add_footer_page_number(slide, 8, total_slides)
    
    # ── SLIDE 9: CONCLUSION ───────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg_and_border(prs, slide)
    
    box = slide.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(10.33), Inches(4.0))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    p = tf.paragraphs[0]
    p.text = "Zusammenfassung"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT_TITLE
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(20)
    
    p2 = tf.add_paragraph()
    p2.text = "Das AALeC-Quizsystem demonstriert die Integration physischer IoT-Controller in eine ereignisgesteuerte Web-Architektur.\n\n" \
              "Der Einsatz von MQTT ermöglicht niedrige Latenzzeiten und lokalen Betrieb ohne externe Internetverbindung.\n\n" \
              "Die funktionale Trennung von Hardware, Backend-Logik und Dashboard sorgt für Wartbarkeit und Erweiterbarkeit."
    p2.alignment = PP_ALIGN.CENTER
    p2.font.name = FONT_BODY
    p2.font.size = Pt(18)
    p2.font.color.rgb = TEXT_MUTED
    p2.space_after = Pt(25)
    
    p3 = tf.add_paragraph()
    p3.text = "Ende der Präsentation"
    p3.alignment = PP_ALIGN.CENTER
    p3.font.name = FONT_TITLE
    p3.font.size = Pt(20)
    p3.font.bold = True
    p3.font.color.rgb = ACCENT_GREEN
    
    add_footer_page_number(slide, 9, total_slides)
    
    # Save Presentation
    output_path = os.path.join(os.path.dirname(__file__), "presentation.pptx")
    prs.save(output_path)
    print(f"Presentation saved successfully at: {output_path}")

if __name__ == "__main__":
    main()
